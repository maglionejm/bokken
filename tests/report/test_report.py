"""Report exports: deterministic, journaled, full coverage, honest appendix."""

from pathlib import Path

import pytest
from pptx import Presentation

from bokken.dossier.model import build_model
from bokken.journal import read_events
from bokken.report.context import build_context, first_sentence
from bokken.report.generate import generate_report, report_exists
from tests.stages.fake_provider import ScriptedProvider
from tests.stages.test_engines_e2e import BRIEF, make_inputs, make_runner


@pytest.fixture(autouse=True)
def home(tmp_path: Path, monkeypatch):
    monkeypatch.setenv("BOKKEN_HOME", str(tmp_path / "home"))
    return tmp_path


@pytest.fixture
def dojo_session(tmp_path: Path) -> Path:
    from bokken.orchestrator import create_session

    brief = {**BRIEF, "inputs": make_inputs(tmp_path)}
    session_dir = create_session(
        "report-dojo",
        brief=brief,
        mode="dojo",
        gate_policy="none",
        config_extra={"panel": {"size": 6, "seed": 11}},
    )
    assert make_runner(session_dir, ScriptedProvider()).run().halt == "completed"
    return session_dir


def deck_text(path: Path) -> str:
    chunks = []
    for slide in Presentation(str(path)).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
    return "\n".join(chunks)


def test_openai_models_have_explicit_list_prices() -> None:
    from bokken.report.context import PRICE_PER_MTOK

    assert PRICE_PER_MTOK["gpt-5"] == (1.25, 10.0)
    assert PRICE_PER_MTOK["gpt-5-mini"] == (0.25, 2.0)
    assert PRICE_PER_MTOK["gpt-5.6-luna"] == (0.2, 1.2)


CACHE_HEAVY_USAGE = {
    "input_tokens": 5_000,
    "output_tokens": 500,
    "cache_read_tokens": 195_000,
    "cache_write_tokens": 10_000,
}


def test_cache_multipliers_are_per_provider() -> None:
    """Anthropic bills cache writes at a premium over input; OpenAI bills none.
    One vendor's multipliers must not be applied to the other's models."""
    from bokken.report.context import call_cost_usd

    # claude-fable-5 lists at $10/$50 per Mtok: 5k input + 500 output +
    # 195k cache reads at 0.1x input + 10k cache writes at 1.25x input.
    assert call_cost_usd("claude-fable-5", CACHE_HEAVY_USAGE) == pytest.approx(0.395)
    # gpt-5 lists at $1.25/$10 per Mtok: cached input at 0.1x, cache writes free.
    assert call_cost_usd("gpt-5", CACHE_HEAVY_USAGE) == pytest.approx(0.035625)
    assert call_cost_usd("gpt-5", {"cache_write_tokens": 10**6}) == 0.0


def test_one_trace_prices_identically_through_cost_and_report_paths(dojo_session: Path) -> None:
    """The `costs` verb and the report appendix are two views of one journal, so
    a cache-heavy call must be quoted at the same number by both."""
    from bokken.journal import Actor, JournalStore
    from bokken.report.context import call_cost_usd, cost_rows

    with JournalStore.open(dojo_session) as store:
        store.append(
            type="model.called",
            stage="empathize",
            actor=Actor(kind="agent", name="model-router"),
            payload={
                "routing_class": "research",
                "model": "claude-fable-5",
                "requested_model": "claude-fable-5",
                "prompt_id": "empathize/interview",
                "prompt_version": "v1",
                "prompt_hash": "h" * 64,
                "request_id": "req-cache",
                "usage": dict(CACHE_HEAVY_USAGE),
                "status": "ok",
            },
        )
    model = build_model(dojo_session)
    rows = cost_rows(model)
    context = build_context(dojo_session, model)
    assert sum(r["cost_usd"] for r in rows) == pytest.approx(context.total_cost_usd, abs=1e-3)

    cached_row = next(r for r in rows if r["model"] == "claude-fable-5")
    cached_line = next(u for u in context.usage if u.model == "claude-fable-5")
    assert cached_row["cache_read"] == cached_line.cache_read_tokens == 195_000
    assert cached_row["cache_write"] == cached_line.cache_write_tokens == 10_000
    assert cached_row["cost_usd"] == pytest.approx(
        call_cost_usd("claude-fable-5", CACHE_HEAVY_USAGE), abs=1e-4
    )
    # The per-model line aggregates the session's other Fable calls too, so it
    # is priced through the same function over its own four buckets.
    assert cached_line.cost_usd == pytest.approx(
        call_cost_usd(
            "claude-fable-5",
            {
                "input_tokens": cached_line.input_tokens,
                "output_tokens": cached_line.output_tokens,
                "cache_read_tokens": cached_line.cache_read_tokens,
                "cache_write_tokens": cached_line.cache_write_tokens,
            },
        )
    )
    assert cached_line.cost_usd > call_cost_usd("claude-fable-5", CACHE_HEAVY_USAGE) * 0.99


def test_export_writes_both_files_journaled_without_model_calls(dojo_session: Path) -> None:
    calls_before = sum(1 for e in read_events(dojo_session) if e.type == "model.called")
    pptx_path, html_path = generate_report(dojo_session)
    assert pptx_path.exists() and html_path.exists()
    kinds = [e.payload["kind"] for e in read_events(dojo_session) if e.type == "artifact.generated"]
    assert "report_deck" in kinds and "report_page" in kinds
    assert report_exists(dojo_session)
    calls_after = sum(1 for e in read_events(dojo_session) if e.type == "model.called")
    assert calls_after == calls_before
    # regeneration is stable
    html_first = html_path.read_text()
    generate_report(dojo_session)
    assert html_path.read_text() == html_first


def test_full_process_coverage_in_html(dojo_session: Path) -> None:
    _, html_path = generate_report(dojo_session)
    html = html_path.read_text()
    model = build_model(dojo_session)
    assert model.problem_statement.resolution in html
    assert "weaker coverage" in html  # a loser with why_lost
    for a in model.assumptions.values():
        assert a.statement in html
        assert (a.score or "untested").upper() in html
    assert model.recommendation.resolution == "iterate"
    assert "iterate" in html
    assert "list-price estimate" in html


def test_honesty_banner_in_both_formats(dojo_session: Path) -> None:
    pptx_path, html_path = generate_report(dojo_session)
    assert "Simulated run." in html_path.read_text()
    text = deck_text(pptx_path)
    assert "SIMULATED RUN" in text
    assert "synthetic" in text


def test_appendix_lists_specs_after_handoff(dojo_session: Path) -> None:
    from bokken.cli import wiring
    from bokken.handoff import finalize_session

    result = finalize_session(dojo_session, lambda store: wiring_router(store))
    assert result.handoff_generated and result.report_generated
    del wiring  # only imported to mirror production call sites
    model = build_model(dojo_session)
    ctx = build_context(dojo_session, model)
    assert ctx.spec_entries, "handoff specs should be summarized"
    for entry in ctx.spec_entries:
        assert entry.sentence.endswith(".")
        assert (dojo_session / entry.path).exists()
    html = (dojo_session / "report" / "report.html").read_text()
    assert ctx.spec_entries[0].path in html


def wiring_router(store):
    from bokken.models import ModelRouter

    return ModelRouter(store, ScriptedProvider())


def test_appendix_surfaces_refusal_for_kill(dojo_session: Path) -> None:
    # Simulate a kill verdict by rebuilding context on a model with a kill recommendation.
    model = build_model(dojo_session)
    model.recommendation.resolution = "kill"
    ctx = build_context(dojo_session, model)
    assert ctx.spec_entries == []
    assert "kill" in (ctx.handoff_refusal or "")
    from bokken.report.page import render_page

    assert "Handoff refused" in render_page(ctx) or "kill" in render_page(ctx)


def test_finalization_is_idempotent_for_report(dojo_session: Path) -> None:
    from bokken.handoff import finalize_session

    first = finalize_session(dojo_session, wiring_router)
    assert first.report_generated
    second = finalize_session(dojo_session, wiring_router)
    assert not second.report_generated
    assert second.summary() == "already finalized"


def test_first_sentence() -> None:
    assert first_sentence("One. Two.") == "One."
    assert first_sentence("No terminal punctuation") == "No terminal punctuation"


# A journal written before the typed read path existed; see
# tests/journal/test_forward_compat.py for what it deliberately contains.
LEGACY_SESSION = Path(__file__).parent.parent / "journal" / "fixtures" / "legacy-session"


def test_legacy_session_still_builds_a_dossier() -> None:
    """The typed readers must produce the same honesty labels on an old journal."""
    model = build_model(LEGACY_SESSION)

    by_class = {e.confidence_class: e for e in model.evidence.values()}
    assert set(by_class) == {"simulated", "reported"}
    assert by_class["simulated"].synthetic is True
    assert by_class["simulated"].speaker == "Marta"
    assert by_class["simulated"].agent is None  # a persona utterance, not an agent's
    assert by_class["simulated"].segment == "commuters"
    assert by_class["simulated"].citations == [{"source_id": "doc-1", "quote": "wrong arrival"}]
    assert by_class["reported"].synthetic is False
    assert by_class["reported"].agent == "founder"

    insight = next(iter(model.insights.values()))
    assert insight.ungrounded is False
    assert insight.synthetic is False  # grounded partly in `reported` evidence

    decision = next(iter(model.decisions.values()))
    assert decision.dissent == [{"actor": "skeptic", "reservation": "the feed is the real problem"}]
    assert decision.requires_real_validation is False
    assert [m.kind for m in model.pivotal_moments] == ["adopted_with_dissent"]

    trace = model.model_traces[0]
    assert trace.usage == {
        "input_tokens": 1200,
        "output_tokens": 340,
        "cache_read_tokens": 800,
        "cache_write_tokens": 0,
    }
    assert [(t.from_stage, t.to_stage, t.loopback) for t in model.transitions] == [
        ("intake", "empathize", False)
    ]
    assert model.abstentions[0].gap == "no operational data in the corpus"


def test_legacy_session_reads_as_a_demo() -> None:
    assert build_context(LEGACY_SESSION, build_model(LEGACY_SESSION)).demo is True


def test_demo_detection_degrades_on_an_unreadable_journal(tmp_path: Path) -> None:
    from bokken.report.context import _is_demo_session

    assert _is_demo_session(tmp_path) is False  # no journal at all
    (tmp_path / "journal.jsonl").write_text("not json\n", encoding="utf-8")
    assert _is_demo_session(tmp_path) is False
