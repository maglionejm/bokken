"""PPTX renderer: house deck grammar (kicker, statement title, blocks, footer)."""

from __future__ import annotations

import contextlib
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from bokken.report.context import ReportContext, split_losers

INK = RGBColor(0x1F, 0x24, 0x30)
GRAY = RGBColor(0x6B, 0x72, 0x80)
ACCENT = RGBColor(0xC7, 0x3E, 0x3A)  # bokken vermillion
ACCENT_DARK = RGBColor(0x8A, 0x2B, 0x28)
PAPER = RGBColor(0xFA, 0xF7, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Helvetica Neue"

SCORE_COLORS = {
    "supported": RGBColor(0x2F, 0x6F, 0x4E),
    "contradicted": ACCENT,
    "untested": GRAY,
}

PAGE_W, PAGE_H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.6)
BODY_W = Inches(12.13)


class Deck:
    def __init__(self, ctx: ReportContext) -> None:
        self.ctx = ctx
        self.prs = Presentation()
        self.prs.slide_width = PAGE_W
        self.prs.slide_height = PAGE_H
        self.page = 0

    # -- primitives ---------------------------------------------------------

    def slide(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])  # blank

    def text(self, slide, x, y, w, h, *, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(x, y, w, h)
        frame = box.text_frame
        frame.word_wrap = True
        frame.paragraphs[0].alignment = align
        return frame

    def run(self, para, text, *, size=10.5, bold=False, color=INK):
        r = para.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        return r

    def para(self, frame, text, *, size=10.5, bold=False, color=INK, before=2, first=False):
        p = frame.paragraphs[0] if first and not frame.paragraphs[0].runs else frame.add_paragraph()
        p.space_before = Pt(0 if first else before)
        self.run(p, text, size=size, bold=bold, color=color)
        return p

    def rect(self, slide, x, y, w, h, fill, *, line=None):
        shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE.RECTANGLE
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if line is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = line
        shape.shadow.inherit = False
        return shape

    def header(self, slide, kicker: str, title: str):
        self.page += 1
        k = self.text(slide, MARGIN, Inches(0.4), BODY_W, Inches(0.3))
        self.para(k, kicker.upper(), size=11, bold=True, color=ACCENT, first=True)
        t = self.text(slide, MARGIN, Inches(0.72), BODY_W, Inches(0.7))
        self.para(t, title, size=26, bold=True, first=True)
        f = self.text(slide, MARGIN, Inches(7.08), BODY_W, Inches(0.3))
        p = self.para(
            f,
            f"Bokken · session {self.ctx.model.name} — generated from the Journal, "
            f"no manual edits        {self.page:02d}",
            size=8,
            color=GRAY,
            first=True,
        )
        p.alignment = PP_ALIGN.LEFT

    def block_title(self, frame, text: str, *, first=True):
        self.para(frame, text, size=12.5, bold=True, color=ACCENT_DARK, first=first)

    def bullets(self, frame, items: list[str], *, size=10.5, first=False):
        for i, item in enumerate(items):
            self.para(frame, f"•  {item}", size=size, first=first and i == 0)

    def chip(self, slide, x, y, w, label: str, body: str, *, fill=INK):
        bar = self.rect(slide, x, y, w, Inches(0.5), fill)
        tf = bar.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self.run(p, label, size=10.5, bold=True, color=WHITE)
        under = self.text(slide, x, y + Inches(0.62), w, Inches(1.9))
        self.para(under, body, size=9.5, first=True)

    def apo(self, slide, stage: str, output_line: str):
        """Agents & activity / Process / Output strip under the title."""
        d = self.ctx.stage_digest.get(stage, {})
        who = (
            ", ".join(list(d.get("personas", []))[:5] + list(d.get("systems", []))) or "facilitator"
        )
        calls = ", ".join(f"{v} {k}" for k, v in sorted(d.get("calls", {}).items())) or "none"
        blocks = [
            (
                "AGENTS & ACTIVITY",
                f"{who}. Calls: {calls} on {', '.join(d.get('models', [])) or '-'}.",
            ),
            ("PROCESS", d.get("process", "")),
            ("OUTPUT", output_line),
        ]
        x, w = MARGIN, Inches(3.95)
        for label, body in blocks:
            frame = self.text(slide, x, Inches(1.5), w, Inches(1.5))
            self.para(frame, label, size=9, bold=True, color=ACCENT_DARK, first=True)
            self.para(frame, body[:330], size=8.5)
            x += w + Inches(0.12)

    APO_BODY_TOP = Inches(3.15)

    def anatomy(self):
        m, c = self.ctx.model, self.ctx
        s = self.slide()
        self.header(s, "how this run worked", "The sequence, the cast, the machinery")
        seq = self.text(s, MARGIN, Inches(1.5), BODY_W, Inches(0.5))
        arrows = "   >   ".join(
            f"{tr.from_stage}->{tr.to_stage}" + (" (loop-back)" if tr.loopback else "")
            for tr in m.transitions
        )
        self.para(seq, arrows[:220], size=10, bold=True, first=True)
        by_panel: dict[str, list] = {}
        for pc in m.personas:
            by_panel.setdefault(pc.panel_kind, []).append(pc)
        left = self.text(s, MARGIN, Inches(2.2), Inches(6.4), Inches(4.4))
        self.block_title(left, "The cast (synthetic, labeled)")
        for kind in ("interview", "ideation", "test"):
            cast = by_panel.get(kind, [])
            if not cast:
                continue
            names = "; ".join(f"{pc.name} ({pc.role})" for pc in cast)
            self.para(left, f"{kind} panel — {names}"[:250], size=9.5, before=6)
        right = self.text(s, Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.4))
        self.block_title(right, "System agents and models")
        self.bullets(
            right,
            [
                "facilitator — stage machinery: programs, clustering, register, verdicts",
                "ui-walker — real-browser walkthrough of the running app",
                "lenses — feasibility vs code, RICE, outcome desirability",
                "skeptic — mandatory on-record challenge",
            ]
            + [f"{u.model}: {u.calls} calls" for u in c.usage],
            size=9.5,
        )

    # -- slides -------------------------------------------------------------

    def cover(self):
        m = self.ctx.model
        s = self.slide()
        self.rect(s, 0, 0, PAGE_W, PAGE_H, PAPER)
        self.rect(s, 0, Inches(6.9), PAGE_W, Inches(0.6), INK)
        f = self.text(s, MARGIN, Inches(2.0), BODY_W, Inches(2.8))
        self.para(
            f, "BOKKEN · DESIGN THINKING RUN REPORT", size=16, bold=True, color=ACCENT, first=True
        )
        headline = self.ctx.headline
        self.para(f, headline, size=40 if len(headline) <= 70 else 28, bold=True, before=10)
        outcome = m.recommendation.resolution if m.recommendation else m.stage
        self.para(
            f,
            f"Session '{m.name}' · mode {m.mode} · status {m.status} · outcome: {outcome}",
            size=14,
            color=GRAY,
            before=10,
        )
        if m.dojo_banner:
            b = self.text(s, MARGIN, Inches(5.5), BODY_W, Inches(1.1))
            self.para(
                b,
                "SIMULATED RUN — produced by an autonomous run against a governed synthetic "
                "persona panel. All persona contributions are simulated and evidence-bounded; "
                "flagged decisions require validation with real users.",
                size=11,
                bold=True,
                first=True,
            )

    def executive_summary(self):
        m, c = self.ctx.model, self.ctx
        s = self.slide()
        self.header(
            s,
            "executive summary",
            f"Outcome: {m.recommendation.resolution if m.recommendation else 'in flight'}",
        )
        left = self.text(s, MARGIN, Inches(1.6), Inches(6.0), Inches(4.4))
        self.block_title(left, "Problem statement")
        self.para(
            left,
            m.problem_statement.resolution if m.problem_statement else "(not reached)",
            size=11,
        )
        self.block_title(left, "Concept advanced", first=False)
        self.para(left, m.concept.resolution if m.concept else "(not reached)", size=11)
        right = self.text(s, Inches(6.8), Inches(1.6), Inches(5.9), Inches(4.4))
        self.block_title(right, "The run in numbers")
        counts = c.register_counts
        self.bullets(
            right,
            [
                f"{len(m.evidence)} evidence items ({c.synthetic_evidence} synthetic, labeled)",
                f"{len(m.insights)} insights · {len(m.options)} options · {len(m.decisions)} decisions on record",
                f"{len(m.assumptions)} assumptions: {counts['supported']} supported · "
                f"{counts['contradicted']} contradicted · {counts['untested']} untested",
                f"{len(m.transitions)} stage transitions ({len(c.loopbacks)} loop-back)",
                f"{sum(u.calls for u in c.usage)} model calls · ~${c.total_cost_usd:,.2f} at list prices",
            ],
            size=11,
        )
        if m.recommendation and m.recommendation.requires_real_validation:
            self.block_title(right, "Validation flag", first=False)
            self.para(
                right,
                "The recommendation rests on simulated evidence and requires real-user validation.",
                size=11,
            )

    def process(self):
        m = self.ctx.model
        s = self.slide()
        self.header(s, "process", "How the run moved — every step journaled")
        stages = ["empathize", "define", "ideate", "prototype", "test", "complete"]
        width, gap = Inches(1.95), Inches(0.08)
        x = MARGIN
        for stage in stages:
            visits = [t for t in m.transitions if t.to_stage == stage]
            body = f"entered {len(visits)}x" if visits else "not reached"
            self.chip(s, x, Inches(1.7), width, stage.upper(), body)
            x += width + gap
        arc = self.text(s, MARGIN, Inches(3.2), BODY_W, Inches(3.4))
        self.block_title(arc, "The arc, in order")
        self.bullets(
            arc,
            [
                f"{t.from_stage} → {t.to_stage}"
                + (" (loop-back)" if t.loopback else "")
                + f": {t.condition}"
                for t in m.transitions
            ],
            size=10,
        )

    def inputs(self):
        m = self.ctx.model
        s = self.slide()
        self.header(s, "inputs", "What the run was grounded in")
        brief_inputs = m.brief.get("inputs", {}) or {}
        left = self.text(s, MARGIN, Inches(1.6), Inches(6.0), Inches(4.9))
        self.block_title(left, "Brief")
        self.bullets(
            left,
            [f"Problem space: {m.brief.get('problem_space', '(none)')}"]
            + [f"Segment: {seg}" for seg in m.brief.get("target_segments", [])]
            + [f"Constraint: {con}" for con in m.brief.get("constraints", [])],
            size=10.5,
        )
        right = self.text(s, Inches(6.8), Inches(1.6), Inches(5.9), Inches(4.9))
        self.block_title(right, "Tangible inputs (typed corpus)")
        lines = []
        if brief_inputs.get("repo"):
            lines.append(f"code · repo: {brief_inputs['repo']}")
        for key, kind in (
            ("documents", "document"),
            ("metrics", "metrics"),
            ("discussions", "discussion"),
        ):
            for item in brief_inputs.get(key, []) or []:
                lines.append(f"{kind}: {item}")
        self.bullets(right, lines or ["(no tangible inputs)"], size=10.5)

    def empathize(self):
        m = self.ctx.model
        s = self.slide()
        self.header(
            s, "intermediate output · empathize", "Evidence, with its confidence class on record"
        )
        self.apo(
            s,
            "empathize",
            f"{len(m.evidence)} evidence items, {len(m.abstentions)} abstention(s), {len(self.ctx.opportunities)} ranked outcomes.",
        )
        grounded = [e for e in m.evidence.values() if e.stage == "empathize"][:5]
        frame = self.text(s, MARGIN, self.APO_BODY_TOP, BODY_W, Inches(3.6))
        self.block_title(frame, f"Sample of {len(m.evidence)} captured items (verbatim)")
        for e in grounded:
            who = e.speaker or e.source
            self.para(frame, f"“{e.content[:220]}”", size=10.5, before=8)
            self.para(
                frame,
                f"— {who} · {e.confidence_class}"
                + (f" · {len(e.citations)} citation(s)" if e.citations else ""),
                size=9,
                color=GRAY,
            )
        if m.abstentions:
            self.para(
                frame,
                f"{len(m.abstentions)} question(s) honestly abstained — carried as research debt "
                "(see negative space).",
                size=10.5,
                bold=True,
                before=10,
            )
        if self.ctx.opportunities:
            self.block_title(frame, "Opportunity ranking (Ulwick)", first=False)
            self.bullets(frame, [o[:180] for o in self.ctx.opportunities[:5]], size=9.5)

    def ui_review(self):
        c = self.ctx
        if not c.ui_review:
            return
        s = self.slide()
        self.header(s, "observed · functional UI review", "The product, exercised first-hand")
        frame = self.text(
            s, MARGIN, Inches(1.6), Inches(7.4) if c.ui_screenshots else BODY_W, Inches(5.2)
        )
        lines = [ln.strip() for ln in c.ui_review.splitlines() if ln.strip()][:16]
        for i, line in enumerate(lines):
            self.para(frame, line.lstrip("#- ")[:150], size=9.5, first=i == 0)
        if c.ui_screenshots:
            shot = c.session_dir / c.ui_screenshots[0]
            if shot.exists():
                # An unreadable image must not kill the export; facts still stand.
                with contextlib.suppress(OSError):
                    s.shapes.add_picture(str(shot), Inches(8.2), Inches(1.6), width=Inches(4.5))

    def define(self):
        m = self.ctx.model
        s = self.slide()
        self.header(
            s, "intermediate output · define", "Problem statement — winner and why the losers lost"
        )
        self.apo(
            s,
            "define",
            "One statement selected; losers preserved with reasons; opportunity coverage in the criteria.",
        )
        frame = self.text(s, MARGIN, self.APO_BODY_TOP, BODY_W, Inches(3.6))
        if not m.problem_statement:
            self.para(frame, "(define was not reached)", first=True)
            return
        self.block_title(frame, "Selected")
        self.para(frame, m.problem_statement.resolution, size=11.5)
        losers = split_losers(m.problem_statement.options)
        if losers:
            self.block_title(frame, "Losers, with the recorded reason", first=False)
            for statement, why in losers[:4]:
                self.para(frame, f"•  {statement[:160]}", size=10)
                self.para(frame, f"    lost: {why[:200]}", size=9.5, color=GRAY)
        reframes = [mv for mv in m.moves if mv.move_id == "hmw_reframe" and mv.executed]
        if reframes:
            self.para(
                frame,
                f"The Kata intervened {len(reframes)}x: solution-shaped statements were reframed "
                "before selection.",
                size=10.5,
                bold=True,
                before=10,
            )

    def ideate(self):
        m = self.ctx.model
        s = self.slide()
        self.header(s, "intermediate output · ideate", "Options on the table, one concept advanced")
        self.apo(
            s,
            "ideate",
            f"{len(m.options)} options with lineage; lens verdicts and dissent on record.",
        )
        frame = self.text(s, MARGIN, self.APO_BODY_TOP, BODY_W, Inches(3.6))
        self.block_title(frame, f"{len(m.options)} options generated")
        if m.concept:
            self.block_title(frame, "Advanced to prototype", first=False)
            self.para(frame, m.concept.resolution, size=11.5)
            for d in m.concept.dissent:
                if isinstance(d, dict):
                    self.para(
                        frame,
                        f"Dissent on record ({d.get('actor', 'panel')}): {d.get('position', '')[:280]}",
                        size=10,
                        color=GRAY,
                        before=8,
                    )
        skeptic = [mv for mv in m.moves if mv.move_id == "skeptic_challenge" and mv.executed]
        if skeptic and skeptic[0].outcome:
            self.block_title(frame, "Skeptic challenge", first=False)
            self.para(frame, skeptic[0].outcome[:300], size=10.5)

    def prototype(self):
        m = self.ctx.model
        s = self.slide()
        self.header(
            s,
            "intermediate output · prototype",
            "Cheapest artifacts against the riskiest assumptions",
        )
        self.apo(
            s,
            "prototype",
            f"{len(m.assumptions)} assumptions; {len(self.ctx.prototype_artifacts)} artifacts hash-journaled.",
        )
        frame = self.text(s, MARGIN, self.APO_BODY_TOP, BODY_W, Inches(3.6))
        fidelity = next(
            (d for d in m.decisions.values() if d.question.startswith("prototype fidelity")), None
        )
        if fidelity:
            self.block_title(frame, "Fidelity decision")
            self.para(frame, fidelity.resolution[:600], size=10)
        real = self.ctx.prototype_artifacts
        if real:
            self.block_title(frame, "Artifacts (hashed, journaled)", first=False)
            self.bullets(
                frame,
                [
                    f"{a.path} · {a.kind} · sha256 {a.content_hash[:12]} · "
                    f"tests {len(a.assumption_ids)} assumption(s)"
                    for a in real
                ],
                size=10.5,
            )

    def concept_research(self):
        mr = self.ctx.market_research
        if not mr:
            return
        s = self.slide()
        self.header(s, "reported · concept research", "The web, on the record")
        left = self.text(s, MARGIN, Inches(1.6), Inches(6.2), Inches(5.0))
        if mr.get("competitors"):
            self.block_title(left, "Competitors and prior art")
            self.bullets(
                left,
                [
                    f"{x.get('name', '')} - {x.get('what', '')} Overlap: {x.get('overlap', '')}"[
                        :200
                    ]
                    for x in mr["competitors"][:5]
                ],
                size=9.5,
            )
        if mr.get("market_signals"):
            self.block_title(left, "Market signals (sourced)", first=False)
            self.bullets(
                left,
                [
                    f"{x.get('stat', '')} ({x.get('source_url', '')})"[:200]
                    for x in mr["market_signals"][:5]
                ],
                size=9.5,
            )
        right = self.text(s, Inches(7.0), Inches(1.6), Inches(5.7), Inches(5.0))
        started = False
        for key, title in (
            ("regulatory", "Regulatory"),
            ("pricing_benchmarks", "Pricing benchmarks"),
            ("differentiation_risks", "Differentiation risks"),
            ("open_questions", "Open questions"),
        ):
            if mr.get(key):
                self.block_title(right, title, first=not started)
                started = True
                self.bullets(right, [str(i)[:160] for i in mr[key][:4]], size=9)

    def test(self):
        m = self.ctx.model
        s = self.slide()
        self.header(
            s, "intermediate output · test", "The assumption register, scored by a fresh panel"
        )
        counts = self.ctx.register_counts
        self.apo(
            s,
            "test",
            f"{counts['supported']} supported, {counts['contradicted']} contradicted, {counts['untested']} untested.",
        )
        y = self.APO_BODY_TOP
        bar_x, bar_w = Inches(9.9), Inches(2.8)
        for a in list(m.assumptions.values())[:8]:
            frame = self.text(s, MARGIN, y, Inches(9.1), Inches(0.5))
            self.para(frame, a.statement[:140], size=9.5, first=True)
            color = SCORE_COLORS[a.score or "untested"]
            self.rect(s, bar_x, y + Inches(0.05), bar_w, Inches(0.28), color)
            label = self.text(s, bar_x, y + Inches(0.05), bar_w, Inches(0.28))
            p = label.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            self.run(p, (a.score or "untested").upper(), size=9, bold=True, color=WHITE)
            y += Inches(0.5)

    def verdict(self):
        m = self.ctx.model
        s = self.slide()
        self.header(
            s,
            "final output · verdict",
            f"Recommendation: {m.recommendation.resolution if m.recommendation else '(none)'}",
        )
        frame = self.text(s, MARGIN, Inches(1.6), BODY_W, Inches(4.9))
        if m.recommendation:
            self.block_title(frame, "Confidence, as recorded")
            note = (
                "Requires real-user validation: the register was scored by a synthetic panel."
                if m.recommendation.requires_real_validation
                else "Scored by humans."
            )
            self.para(frame, note, size=10.5)
        loop = next(
            (mv for mv in m.moves if mv.move_id == "loopback_proposal" and mv.executed), None
        )
        if loop and loop.outcome:
            self.block_title(frame, "Loop-back proposal fired at test", first=False)
            self.para(frame, loop.outcome[:600], size=10.5)
        for moment in m.pivotal_moments:
            self.para(frame, f"•  {moment.description[:260]}", size=10, before=6)

    def negative_space(self):
        m = self.ctx.model
        debt = m.negative_space.research_debt
        if not debt:
            return
        s = self.slide()
        self.header(s, "negative space", "What this run did not do — open research debt")
        frame = self.text(s, MARGIN, Inches(1.6), BODY_W, Inches(4.9))
        seen: set[str] = set()
        shown = 0
        for item in debt:
            if item.question in seen or shown >= 7:
                continue
            seen.add(item.question)
            shown += 1
            self.para(frame, f"•  {item.question[:220]}", size=10.5, first=shown == 1)
        self.para(
            frame,
            "These are real-user questions the synthetic panel refused to answer from the corpus.",
            size=10,
            color=GRAY,
            before=10,
        )

    def model_ops(self):
        c = self.ctx
        s = self.slide()
        self.header(
            s,
            "final output · model operations",
            "Every call journaled; cost estimated at list prices",
        )
        frame = self.text(s, MARGIN, Inches(1.6), BODY_W, Inches(4.4))
        self.block_title(frame, "Usage by model")
        self.bullets(
            frame,
            [
                f"{u.model}: {u.calls} calls · {u.input_tokens:,} in / {u.output_tokens:,} out "
                f"· ~${u.cost_usd:,.2f}"
                for u in c.usage
            ]
            + [f"Total: ~${c.total_cost_usd:,.2f} (list-price estimate)"],
            size=11,
        )
        if c.dossier_paths:
            self.block_title(frame, "Companion records", first=False)
            records = [*c.dossier_paths, "journal.jsonl (hash-chained, append-only)"]
            self.bullets(frame, records, size=10.5)

    def appendix(self):
        c = self.ctx
        s = self.slide()
        self.header(s, "appendix", "Specifications handed off — one line each")
        frame = self.text(s, MARGIN, Inches(1.6), BODY_W, Inches(4.9))
        if c.spec_entries:
            for i, entry in enumerate(c.spec_entries):
                self.para(frame, f"{entry.capability} — {entry.sentence}", size=10.5, first=i == 0)
                self.para(frame, f"full spec: {entry.path}", size=9, color=GRAY)
        else:
            self.block_title(frame, "No handoff package")
            self.para(
                frame,
                f"Handoff refused: {c.handoff_refusal or 'not generated yet'}.",
                size=11,
            )

    def render(self, out_path: Path) -> None:
        self.cover()
        self.executive_summary()
        self.anatomy()
        self.process()
        self.inputs()
        self.empathize()
        self.ui_review()
        self.define()
        self.ideate()
        self.prototype()
        self.concept_research()
        self.test()
        self.verdict()
        self.negative_space()
        self.model_ops()
        self.appendix()
        out_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(out_path))


def render_deck(ctx: ReportContext, out_path: Path) -> None:
    Deck(ctx).render(out_path)
